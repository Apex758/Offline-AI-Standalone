"""
File parser module for the File Explorer feature.
Extracts text content from various document formats.
"""

import io
import csv
from pathlib import Path


def parse_docx(content: bytes) -> dict:
    """Extract text from a Word document."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = '\n'.join(paragraphs)

        # Extract metadata
        core = doc.core_properties
        metadata = {}
        if core.title:
            metadata['title'] = core.title
        if core.author:
            metadata['author'] = core.author

        return {'text': text, 'paragraphs': paragraphs, 'metadata': metadata}
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_pptx(content: bytes) -> dict:
    """Extract text from a PowerPoint presentation."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        all_text = []

        for slide in prs.slides:
            slide_data = {'title': '', 'bullets': []}
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        # First text found is likely the title
                        if not slide_data['title']:
                            slide_data['title'] = text
                        else:
                            slide_data['bullets'].append(text)
                        all_text.append(text)
            slides.append(slide_data)

        return {
            'text': '\n'.join(all_text),
            'slides': slides,
            'metadata': {'slide_count': len(slides)}
        }
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_xlsx(content: bytes) -> dict:
    """Extract data from an Excel spreadsheet."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets = []
        all_text = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            headers = []
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                str_row = [str(cell) if cell is not None else '' for cell in row]
                if row_idx == 0:
                    headers = str_row
                rows.append(str_row)
                all_text.extend([c for c in str_row if c])

            sheets.append({
                'name': sheet_name,
                'headers': headers,
                'rows': rows,
                'row_count': len(rows)
            })

        wb.close()
        return {
            'text': '\n'.join(all_text),
            'sheets': sheets,
            'metadata': {'sheet_count': len(sheets)}
        }
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_pdf(content: bytes) -> dict:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = []
        all_text = []

        for page in reader.pages:
            page_text = page.extract_text() or ''
            pages.append(page_text)
            all_text.append(page_text)

        return {
            'text': '\n'.join(all_text),
            'pages': pages,
            'metadata': {'page_count': len(pages)}
        }
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_text(content: bytes) -> dict:
    """Parse plain text, markdown, or CSV files."""
    try:
        text = content.decode('utf-8', errors='replace')
        return {'text': text}
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_csv_file(content: bytes) -> dict:
    """Parse a CSV file into structured data."""
    try:
        text = content.decode('utf-8', errors='replace')
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        headers = rows[0] if rows else []
        return {
            'text': text,
            'headers': headers,
            'rows': rows,
            'metadata': {'row_count': len(rows)}
        }
    except Exception as e:
        return {'text': '', 'error': str(e)}


def parse_file(content: bytes, filename: str) -> dict:
    """Route to the correct parser based on file extension."""
    ext = Path(filename).suffix.lower()

    parsers = {
        '.docx': parse_docx,
        '.doc': parse_docx,
        '.pptx': parse_pptx,
        '.ppt': parse_pptx,
        '.xlsx': parse_xlsx,
        '.xls': parse_xlsx,
        '.pdf': parse_pdf,
        '.csv': parse_csv_file,
        '.txt': parse_text,
        '.md': parse_text,
    }

    parser = parsers.get(ext, parse_text)
    result = parser(content)
    result['filename'] = filename
    result['extension'] = ext
    return result


def get_preview(content: bytes, filename: str, max_words: int = 500) -> dict:
    """Get a truncated preview of file content."""
    result = parse_file(content, filename)
    text = result.get('text', '')

    # Truncate to max_words
    words = text.split()
    if len(words) > max_words:
        result['text'] = ' '.join(words[:max_words]) + '...'
        result['truncated'] = True
    else:
        result['truncated'] = False

    return result
